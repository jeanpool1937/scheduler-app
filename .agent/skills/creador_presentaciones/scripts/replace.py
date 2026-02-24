#!/usr/bin/env python3
"""
Aplicar reemplazos de texto a presentación PowerPoint.

Uso:
    python replace.py <input.pptx> <replacements.json> <output.pptx>

El JSON de reemplazos debe tener la estructura producida por inventory.py.
TODAS las formas de texto identificadas por inventory.py tendrán su texto limpiado
a menos que se especifique "paragraphs" en los reemplazos para esa forma.
"""

import json
import sys
from pathlib import Path
from typing import Any, Dict, List

from inventory import InventoryData, extract_text_inventory
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


def clear_paragraph_bullets(paragraph):
    """Limpiar formato de viñeta de un párrafo."""
    pPr = paragraph._element.get_or_add_pPr()

    for child in list(pPr):
        if (
            child.tag.endswith("buChar")
            or child.tag.endswith("buNone")
            or child.tag.endswith("buAutoNum")
            or child.tag.endswith("buFont")
        ):
            pPr.remove(child)

    return pPr


def apply_paragraph_properties(paragraph, para_data: Dict[str, Any]):
    """Aplicar propiedades de formato a un párrafo."""
    text = para_data.get("text", "")
    pPr = clear_paragraph_bullets(paragraph)

    # Manejar formato de viñeta
    if para_data.get("bullet", False):
        level = para_data.get("level", 0)
        paragraph.level = level

        font_size = para_data.get("font_size", 18.0)
        level_indent_emu = int((font_size * (1.6 + level * 1.6)) * 12700)
        hanging_indent_emu = int(-font_size * 0.8 * 12700)

        pPr.attrib["marL"] = str(level_indent_emu)
        pPr.attrib["indent"] = str(hanging_indent_emu)

        buChar = OxmlElement("a:buChar")
        buChar.set("char", "•")
        pPr.append(buChar)

        if "alignment" not in para_data:
            paragraph.alignment = PP_ALIGN.LEFT
    else:
        pPr.attrib["marL"] = "0"
        pPr.attrib["indent"] = "0"

        buNone = OxmlElement("a:buNone")
        pPr.insert(0, buNone)

    # Aplicar alineación
    if "alignment" in para_data:
        alignment_map = {
            "LEFT": PP_ALIGN.LEFT,
            "CENTER": PP_ALIGN.CENTER,
            "RIGHT": PP_ALIGN.RIGHT,
            "JUSTIFY": PP_ALIGN.JUSTIFY,
        }
        if para_data["alignment"] in alignment_map:
            paragraph.alignment = alignment_map[para_data["alignment"]]

    # Aplicar espaciado
    if "space_before" in para_data:
        paragraph.space_before = Pt(para_data["space_before"])
    if "space_after" in para_data:
        paragraph.space_after = Pt(para_data["space_after"])
    if "line_spacing" in para_data:
        paragraph.line_spacing = Pt(para_data["line_spacing"])

    # Aplicar formato a nivel de run
    if not paragraph.runs:
        run = paragraph.add_run()
        run.text = text
    else:
        run = paragraph.runs[0]
        run.text = text

    apply_font_properties(run, para_data)


def apply_font_properties(run, para_data: Dict[str, Any]):
    """Aplicar propiedades de fuente a un run de texto."""
    if "bold" in para_data:
        run.font.bold = para_data["bold"]
    if "italic" in para_data:
        run.font.italic = para_data["italic"]
    if "underline" in para_data:
        run.font.underline = para_data["underline"]
    if "font_size" in para_data:
        run.font.size = Pt(para_data["font_size"])
    if "font_name" in para_data:
        run.font.name = para_data["font_name"]

    # Aplicar color - preferir RGB, recurrir a theme_color
    if "color" in para_data:
        color_hex = para_data["color"].lstrip("#")
        if len(color_hex) == 6:
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            run.font.color.rgb = RGBColor(r, g, b)
    elif "theme_color" in para_data:
        theme_name = para_data["theme_color"]
        try:
            run.font.color.theme_color = getattr(MSO_THEME_COLOR, theme_name)
        except AttributeError:
            print(f"  ADVERTENCIA: Nombre de color de tema desconocido '{theme_name}'")


def detect_frame_overflow(inventory: InventoryData) -> Dict[str, Dict[str, float]]:
    """Detectar overflow de texto en formas."""
    overflow_map = {}

    for slide_key, shapes_dict in inventory.items():
        for shape_key, shape_data in shapes_dict.items():
            if shape_data.frame_overflow_bottom is not None:
                if slide_key not in overflow_map:
                    overflow_map[slide_key] = {}
                overflow_map[slide_key][shape_key] = shape_data.frame_overflow_bottom

    return overflow_map


def validate_replacements(inventory: InventoryData, replacements: Dict) -> List[str]:
    """Validar que todas las formas en reemplazos existen en el inventario."""
    errors = []

    for slide_key, shapes_data in replacements.items():
        if not slide_key.startswith("slide-"):
            continue

        if slide_key not in inventory:
            errors.append(f"Diapositiva '{slide_key}' no encontrada en inventario")
            continue

        for shape_key in shapes_data.keys():
            if shape_key not in inventory[slide_key]:
                unused_with_content = []
                for k in inventory[slide_key].keys():
                    if k not in shapes_data:
                        shape_data = inventory[slide_key][k]
                        paragraphs = shape_data.paragraphs
                        if paragraphs and paragraphs[0].text:
                            first_text = paragraphs[0].text[:50]
                            if len(paragraphs[0].text) > 50:
                                first_text += "..."
                            unused_with_content.append(f"{k} ('{first_text}')")
                        else:
                            unused_with_content.append(k)

                errors.append(
                    f"Forma '{shape_key}' no encontrada en '{slide_key}'. "
                    f"Formas sin reemplazos: {', '.join(sorted(unused_with_content)) if unused_with_content else 'ninguna'}"
                )

    return errors


def check_duplicate_keys(pairs):
    """Verificar claves duplicadas al cargar JSON."""
    result = {}
    for key, value in pairs:
        if key in result:
            raise ValueError(f"Clave duplicada encontrada en JSON: '{key}'")
        result[key] = value
    return result


def apply_replacements(pptx_file: str, json_file: str, output_file: str):
    """Aplicar reemplazos de texto desde JSON a presentación PowerPoint."""

    prs = Presentation(pptx_file)
    inventory = extract_text_inventory(Path(pptx_file), prs)
    original_overflow = detect_frame_overflow(inventory)

    with open(json_file, "r", encoding="utf-8") as f:
        replacements = json.load(f, object_pairs_hook=check_duplicate_keys)

    errors = validate_replacements(inventory, replacements)
    if errors:
        print("ERROR: Formas inválidas en JSON de reemplazo:")
        for error in errors:
            print(f"  - {error}")
        print("\nPor favor verifica el inventario y actualiza tu JSON de reemplazo.")
        raise ValueError(f"Encontrados {len(errors)} error(es) de validación")

    shapes_processed = 0
    shapes_cleared = 0
    shapes_replaced = 0

    for slide_key, shapes_dict in inventory.items():
        if not slide_key.startswith("slide-"):
            continue

        slide_index = int(slide_key.split("-")[1])

        if slide_index >= len(prs.slides):
            print(f"Advertencia: Diapositiva {slide_index} no encontrada")
            continue

        for shape_key, shape_data in shapes_dict.items():
            shapes_processed += 1

            shape = shape_data.shape
            if not shape:
                print(f"Advertencia: {shape_key} no tiene referencia de forma")
                continue

            text_frame = shape.text_frame
            text_frame.clear()
            shapes_cleared += 1

            replacement_shape_data = replacements.get(slide_key, {}).get(shape_key, {})
            if "paragraphs" not in replacement_shape_data:
                continue

            shapes_replaced += 1

            for i, para_data in enumerate(replacement_shape_data["paragraphs"]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                apply_paragraph_properties(p, para_data)

    # Verificar problemas después de reemplazos
    import tempfile

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = Path(tmp.name)
        prs.save(str(tmp_path))

    try:
        updated_inventory = extract_text_inventory(tmp_path)
        updated_overflow = detect_frame_overflow(updated_inventory)
    finally:
        tmp_path.unlink()

    overflow_errors = []
    for slide_key, shape_overflows in updated_overflow.items():
        for shape_key, new_overflow in shape_overflows.items():
            original = original_overflow.get(slide_key, {}).get(shape_key, 0.0)

            if new_overflow > original + 0.01:
                increase = new_overflow - original
                overflow_errors.append(
                    f'{slide_key}/{shape_key}: overflow empeoró en {increase:.2f}" '
                    f'(era {original:.2f}", ahora {new_overflow:.2f}")'
                )

    warnings = []
    for slide_key, shapes_dict in updated_inventory.items():
        for shape_key, shape_data in shapes_dict.items():
            if shape_data.warnings:
                for warning in shape_data.warnings:
                    warnings.append(f"{slide_key}/{shape_key}: {warning}")

    if overflow_errors or warnings:
        print("\nERROR: Problemas detectados en salida de reemplazo:")
        if overflow_errors:
            print("\nOverflow de texto empeoró:")
            for error in overflow_errors:
                print(f"  - {error}")
        if warnings:
            print("\nAdvertencias de formato:")
            for warning in warnings:
                print(f"  - {warning}")
        print("\nPor favor corrige estos problemas antes de guardar.")
        raise ValueError(
            f"Encontrados {len(overflow_errors)} error(es) de overflow y {len(warnings)} advertencia(s)"
        )

    prs.save(output_file)

    print(f"Guardada presentación actualizada en: {output_file}")
    print(f"Procesadas {len(prs.slides)} diapositivas")
    print(f"  - Formas procesadas: {shapes_processed}")
    print(f"  - Formas limpiadas: {shapes_cleared}")
    print(f"  - Formas reemplazadas: {shapes_replaced}")


def main():
    """Punto de entrada principal para uso en línea de comandos."""
    if len(sys.argv) != 4:
        print(__doc__)
        sys.exit(1)

    input_pptx = Path(sys.argv[1])
    replacements_json = Path(sys.argv[2])
    output_pptx = Path(sys.argv[3])

    if not input_pptx.exists():
        print(f"Error: Archivo de entrada '{input_pptx}' no encontrado")
        sys.exit(1)

    if not replacements_json.exists():
        print(f"Error: Archivo JSON de reemplazos '{replacements_json}' no encontrado")
        sys.exit(1)

    try:
        apply_replacements(str(input_pptx), str(replacements_json), str(output_pptx))
    except Exception as e:
        print(f"Error aplicando reemplazos: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
