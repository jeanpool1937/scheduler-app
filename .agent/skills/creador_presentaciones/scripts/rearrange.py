#!/usr/bin/env python3
"""
Reordenar diapositivas de PowerPoint basándose en una secuencia de índices.

Uso:
    python rearrange.py template.pptx output.pptx 0,34,34,50,52

Esto creará output.pptx usando diapositivas de template.pptx en el orden especificado.
Las diapositivas pueden repetirse (ej: 34 aparece dos veces).
"""

import argparse
import shutil
import sys
from copy import deepcopy
from pathlib import Path

import six
from pptx import Presentation


def main():
    parser = argparse.ArgumentParser(
        description="Reordenar diapositivas de PowerPoint basándose en una secuencia de índices.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Ejemplos:
  python rearrange.py template.pptx output.pptx 0,34,34,50,52
    Crea output.pptx usando diapositivas 0, 34 (dos veces), 50 y 52 de template.pptx

  python rearrange.py template.pptx output.pptx 5,3,1,2,4
    Crea output.pptx con diapositivas reordenadas según se especifica

Nota: Los índices de diapositiva son base-0 (primera diapositiva es 0, segunda es 1, etc.)
        """,
    )

    parser.add_argument("template", help="Ruta al archivo PPTX de plantilla")
    parser.add_argument("output", help="Ruta para el archivo PPTX de salida")
    parser.add_argument(
        "sequence", help="Secuencia de índices de diapositiva separados por comas (base-0)"
    )

    args = parser.parse_args()

    # Parsear la secuencia de diapositivas
    try:
        slide_sequence = [int(x.strip()) for x in args.sequence.split(",")]
    except ValueError:
        print(
            "Error: Formato de secuencia inválido. Usa enteros separados por comas (ej: 0,34,34,50,52)"
        )
        sys.exit(1)

    # Verificar que la plantilla existe
    template_path = Path(args.template)
    if not template_path.exists():
        print(f"Error: Archivo de plantilla no encontrado: {args.template}")
        sys.exit(1)

    # Crear directorio de salida si es necesario
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        rearrange_presentation(template_path, output_path, slide_sequence)
    except ValueError as e:
        print(f"Error: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"Error procesando presentación: {e}")
        sys.exit(1)


def duplicate_slide(pres, index):
    """Duplicar una diapositiva en la presentación."""
    source = pres.slides[index]

    # Usar el layout de la diapositiva fuente para preservar formato
    new_slide = pres.slides.add_slide(source.slide_layout)

    # Recopilar todas las relaciones de imagen y multimedia de la diapositiva fuente
    image_rels = {}
    for rel_id, rel in six.iteritems(source.part.rels):
        if "image" in rel.reltype or "media" in rel.reltype:
            image_rels[rel_id] = rel

    # CRÍTICO: Limpiar formas placeholder para evitar duplicados
    for shape in new_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)

    # Copiar todas las formas de la fuente
    for shape in source.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        # Manejar formas de imagen - necesita actualizar la referencia blip
        blips = new_el.xpath(".//a:blip[@r:embed]")
        for blip in blips:
            old_rId = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if old_rId in image_rels:
                old_rel = image_rels[old_rId]
                new_rId = new_slide.part.rels.get_or_add(
                    old_rel.reltype, old_rel._target
                )
                blip.set(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                    new_rId,
                )

    # Copiar relaciones adicionales de imagen/multimedia
    for rel_id, rel in image_rels.items():
        try:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)
        except Exception:
            pass

    return new_slide


def delete_slide(pres, index):
    """Eliminar una diapositiva de la presentación."""
    rId = pres.slides._sldIdLst[index].rId
    pres.part.drop_rel(rId)
    del pres.slides._sldIdLst[index]


def reorder_slides(pres, slide_index, target_index):
    """Mover una diapositiva de una posición a otra."""
    slides = pres.slides._sldIdLst

    # Remover elemento de diapositiva de la posición actual
    slide_element = slides[slide_index]
    slides.remove(slide_element)

    # Insertar en posición destino
    slides.insert(target_index, slide_element)


def rearrange_presentation(template_path, output_path, slide_sequence):
    """
    Crear una nueva presentación con diapositivas de la plantilla en el orden especificado.

    Args:
        template_path: Ruta al archivo PPTX de plantilla
        output_path: Ruta para el archivo PPTX de salida
        slide_sequence: Lista de índices de diapositiva (base-0) a incluir
    """
    # Copiar plantilla para preservar dimensiones y tema
    if template_path != output_path:
        shutil.copy2(template_path, output_path)
        prs = Presentation(output_path)
    else:
        prs = Presentation(template_path)

    total_slides = len(prs.slides)

    # Validar índices
    for idx in slide_sequence:
        if idx < 0 or idx >= total_slides:
            raise ValueError(f"Índice de diapositiva {idx} fuera de rango (0-{total_slides - 1})")

    # Rastrear diapositivas originales y sus duplicados
    slide_map = []
    duplicated = {}

    # Paso 1: DUPLICAR diapositivas repetidas
    print(f"Procesando {len(slide_sequence)} diapositivas de la plantilla...")
    for i, template_idx in enumerate(slide_sequence):
        if template_idx in duplicated and duplicated[template_idx]:
            slide_map.append(duplicated[template_idx].pop(0))
            print(f"  [{i}] Usando duplicado de diapositiva {template_idx}")
        elif slide_sequence.count(template_idx) > 1 and template_idx not in duplicated:
            slide_map.append(template_idx)
            duplicates = []
            count = slide_sequence.count(template_idx) - 1
            print(
                f"  [{i}] Usando diapositiva original {template_idx}, creando {count} duplicado(s)"
            )
            for _ in range(count):
                duplicate_slide(prs, template_idx)
                duplicates.append(len(prs.slides) - 1)
            duplicated[template_idx] = duplicates
        else:
            slide_map.append(template_idx)
            print(f"  [{i}] Usando diapositiva original {template_idx}")

    # Paso 2: ELIMINAR diapositivas no deseadas (trabajar hacia atrás)
    slides_to_keep = set(slide_map)
    print(f"\nEliminando {len(prs.slides) - len(slides_to_keep)} diapositivas no usadas...")
    for i in range(len(prs.slides) - 1, -1, -1):
        if i not in slides_to_keep:
            delete_slide(prs, i)
            slide_map = [idx - 1 if idx > i else idx for idx in slide_map]

    # Paso 3: REORDENAR a secuencia final
    print(f"Reordenando {len(slide_map)} diapositivas a secuencia final...")
    for target_pos in range(len(slide_map)):
        current_pos = slide_map[target_pos]
        if current_pos != target_pos:
            reorder_slides(prs, current_pos, target_pos)
            for i in range(len(slide_map)):
                if slide_map[i] > current_pos and slide_map[i] <= target_pos:
                    slide_map[i] -= 1
                elif slide_map[i] < current_pos and slide_map[i] >= target_pos:
                    slide_map[i] += 1
            slide_map[target_pos] = target_pos

    # Guardar la presentación
    prs.save(output_path)
    print(f"\nGuardada presentación reordenada en: {output_path}")
    print(f"La presentación final tiene {len(prs.slides)} diapositivas")


if __name__ == "__main__":
    main()
