#!/usr/bin/env python3
"""
Crear cuadrículas de miniaturas de diapositivas de PowerPoint.

Crea un layout de cuadrícula de miniaturas de diapositivas con columnas configurables (máx 6).
Cada cuadrícula contiene hasta cols×(cols+1) imágenes. Para presentaciones con más
diapositivas, se crean múltiples archivos de cuadrícula numerados automáticamente.

Salida:
- Cuadrícula única: {prefijo}.jpg (si las diapositivas caben en una cuadrícula)
- Múltiples cuadrículas: {prefijo}-1.jpg, {prefijo}-2.jpg, etc.

Límites de cuadrícula por número de columnas:
- 3 cols: máx 12 diapositivas por cuadrícula (3×4)
- 4 cols: máx 20 diapositivas por cuadrícula (4×5)
- 5 cols: máx 30 diapositivas por cuadrícula (5×6) [por defecto]
- 6 cols: máx 42 diapositivas por cuadrícula (6×7)

Uso:
    python thumbnail.py input.pptx [prefijo_salida] [--cols N] [--outline-placeholders]
"""

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

# Constantes
THUMBNAIL_WIDTH = 300  # Ancho fijo de miniatura en píxeles
CONVERSION_DPI = 100  # DPI para conversión PDF a imagen
MAX_COLS = 6  # Número máximo de columnas
DEFAULT_COLS = 5  # Número de columnas por defecto
JPEG_QUALITY = 95  # Calidad de compresión JPEG

# Constantes de layout de cuadrícula
GRID_PADDING = 20  # Espaciado entre miniaturas
BORDER_WIDTH = 2  # Ancho de borde alrededor de miniaturas
FONT_SIZE_RATIO = 0.12  # Tamaño de fuente como fracción del ancho de miniatura
LABEL_PADDING_RATIO = 0.4  # Espaciado de etiqueta como fracción del tamaño de fuente


def main():
    parser = argparse.ArgumentParser(
        description="Crear cuadrículas de miniaturas de diapositivas PowerPoint."
    )
    parser.add_argument("input", help="Archivo PowerPoint de entrada (.pptx)")
    parser.add_argument(
        "output_prefix",
        nargs="?",
        default="thumbnails",
        help="Prefijo de salida para archivos de imagen (por defecto: thumbnails)",
    )
    parser.add_argument(
        "--cols",
        type=int,
        default=DEFAULT_COLS,
        help=f"Número de columnas (por defecto: {DEFAULT_COLS}, máx: {MAX_COLS})",
    )
    parser.add_argument(
        "--outline-placeholders",
        action="store_true",
        help="Resaltar placeholders de texto con borde de color",
    )

    args = parser.parse_args()

    # Validar columnas
    cols = min(args.cols, MAX_COLS)
    if args.cols > MAX_COLS:
        print(f"Advertencia: Columnas limitadas a {MAX_COLS} (solicitado {args.cols})")

    # Validar entrada
    input_path = Path(args.input)
    if not input_path.exists() or input_path.suffix.lower() != ".pptx":
        print(f"Error: Archivo PowerPoint inválido: {args.input}")
        sys.exit(1)

    # Construir ruta de salida (siempre JPG)
    output_path = Path(f"{args.output_prefix}.jpg")

    print(f"Procesando: {args.input}")

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            # Obtener regiones de placeholder si el resaltado está habilitado
            placeholder_regions = None
            slide_dimensions = None
            if args.outline_placeholders:
                print("Extrayendo regiones de placeholder...")
                placeholder_regions, slide_dimensions = get_placeholder_regions(
                    input_path
                )
                if placeholder_regions:
                    print(f"Encontrados placeholders en {len(placeholder_regions)} diapositivas")

            # Convertir diapositivas a imágenes
            slide_images = convert_to_images(input_path, Path(temp_dir), CONVERSION_DPI)
            if not slide_images:
                print("Error: No se encontraron diapositivas")
                sys.exit(1)

            print(f"Encontradas {len(slide_images)} diapositivas")

            # Crear cuadrículas (máx cols×(cols+1) imágenes por cuadrícula)
            grid_files = create_grids(
                slide_images,
                cols,
                THUMBNAIL_WIDTH,
                output_path,
                placeholder_regions,
                slide_dimensions,
            )

            # Imprimir archivos guardados
            print(f"Creadas {len(grid_files)} cuadrícula(s):")
            for grid_file in grid_files:
                print(f"  - {grid_file}")

    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


def create_hidden_slide_placeholder(size):
    """Crear imagen placeholder para diapositivas ocultas."""
    img = Image.new("RGB", size, color="#F0F0F0")
    draw = ImageDraw.Draw(img)
    line_width = max(5, min(size) // 100)
    draw.line([(0, 0), size], fill="#CCCCCC", width=line_width)
    draw.line([(size[0], 0), (0, size[1])], fill="#CCCCCC", width=line_width)
    return img


def get_placeholder_regions(pptx_path):
    """Extraer TODAS las regiones de texto de la presentación."""
    try:
        from inventory import extract_text_inventory
    except ImportError:
        print("Advertencia: No se pudo importar inventory.py, omitiendo placeholders")
        return None, None
        
    prs = Presentation(str(pptx_path))
    inventory = extract_text_inventory(pptx_path, prs)
    placeholder_regions = {}

    # Obtener dimensiones de diapositiva en pulgadas
    slide_width_inches = (prs.slide_width or 9144000) / 914400.0
    slide_height_inches = (prs.slide_height or 5143500) / 914400.0

    for slide_key, shapes in inventory.items():
        slide_idx = int(slide_key.split("-")[1])
        regions = []

        for shape_key, shape_data in shapes.items():
            regions.append(
                {
                    "left": shape_data.left,
                    "top": shape_data.top,
                    "width": shape_data.width,
                    "height": shape_data.height,
                }
            )

        if regions:
            placeholder_regions[slide_idx] = regions

    return placeholder_regions, (slide_width_inches, slide_height_inches)


def convert_to_images(pptx_path, temp_dir, dpi):
    """Convertir PowerPoint a imágenes vía PDF, manejando diapositivas ocultas."""
    print("Analizando presentación...")
    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)

    # Encontrar diapositivas ocultas (indexación base-1 para mostrar)
    hidden_slides = {
        idx + 1
        for idx, slide in enumerate(prs.slides)
        if slide.element.get("show") == "0"
    }

    print(f"Total de diapositivas: {total_slides}")
    if hidden_slides:
        print(f"Diapositivas ocultas: {sorted(hidden_slides)}")

    pdf_path = temp_dir / f"{pptx_path.stem}.pdf"

    # Convertir a PDF
    print("Convirtiendo a PDF...")
    result = subprocess.run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(temp_dir),
            str(pptx_path),
        ],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0 or not pdf_path.exists():
        raise RuntimeError("Falló la conversión a PDF")

    # Convertir PDF a imágenes
    print(f"Convirtiendo a imágenes a {dpi} DPI...")
    result = subprocess.run(
        ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf_path), str(temp_dir / "slide")],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        raise RuntimeError("Falló la conversión a imágenes")

    visible_images = sorted(temp_dir.glob("slide-*.jpg"))

    # Crear lista completa con placeholders para diapositivas ocultas
    all_images = []
    visible_idx = 0

    # Obtener dimensiones de placeholder de la primera diapositiva visible
    if visible_images:
        with Image.open(visible_images[0]) as img:
            placeholder_size = img.size
    else:
        placeholder_size = (1920, 1080)

    for slide_num in range(1, total_slides + 1):
        if slide_num in hidden_slides:
            placeholder_path = temp_dir / f"hidden-{slide_num:03d}.jpg"
            placeholder_img = create_hidden_slide_placeholder(placeholder_size)
            placeholder_img.save(placeholder_path, "JPEG")
            all_images.append(placeholder_path)
        else:
            if visible_idx < len(visible_images):
                all_images.append(visible_images[visible_idx])
                visible_idx += 1

    return all_images


def create_grids(
    image_paths,
    cols,
    width,
    output_path,
    placeholder_regions=None,
    slide_dimensions=None,
):
    """Crear múltiples cuadrículas de miniaturas, máx cols×(cols+1) imágenes por cuadrícula."""
    max_images_per_grid = cols * (cols + 1)
    grid_files = []

    print(
        f"Creando cuadrículas con {cols} columnas (máx {max_images_per_grid} imágenes por cuadrícula)"
    )

    for chunk_idx, start_idx in enumerate(
        range(0, len(image_paths), max_images_per_grid)
    ):
        end_idx = min(start_idx + max_images_per_grid, len(image_paths))
        chunk_images = image_paths[start_idx:end_idx]

        grid = create_grid(
            chunk_images, cols, width, start_idx, placeholder_regions, slide_dimensions
        )

        if len(image_paths) <= max_images_per_grid:
            grid_filename = output_path
        else:
            stem = output_path.stem
            suffix = output_path.suffix
            grid_filename = output_path.parent / f"{stem}-{chunk_idx + 1}{suffix}"

        grid_filename.parent.mkdir(parents=True, exist_ok=True)
        grid.save(str(grid_filename), quality=JPEG_QUALITY)
        grid_files.append(str(grid_filename))

    return grid_files


def create_grid(
    image_paths,
    cols,
    width,
    start_slide_num=0,
    placeholder_regions=None,
    slide_dimensions=None,
):
    """Crear cuadrícula de miniaturas de imágenes de diapositivas."""
    font_size = int(width * FONT_SIZE_RATIO)
    label_padding = int(font_size * LABEL_PADDING_RATIO)

    with Image.open(image_paths[0]) as img:
        aspect = img.height / img.width
    height = int(width * aspect)

    rows = (len(image_paths) + cols - 1) // cols
    grid_w = cols * width + (cols + 1) * GRID_PADDING
    grid_h = rows * (height + font_size + label_padding * 2) + (rows + 1) * GRID_PADDING

    grid = Image.new("RGB", (grid_w, grid_h), "white")
    draw = ImageDraw.Draw(grid)

    try:
        font = ImageFont.load_default(size=font_size)
    except Exception:
        font = ImageFont.load_default()

    for i, img_path in enumerate(image_paths):
        row, col = i // cols, i % cols
        x = col * width + (col + 1) * GRID_PADDING
        y_base = (
            row * (height + font_size + label_padding * 2) + (row + 1) * GRID_PADDING
        )

        label = f"{start_slide_num + i}"
        bbox = draw.textbbox((0, 0), label, font=font)
        text_w = bbox[2] - bbox[0]
        draw.text(
            (x + (width - text_w) // 2, y_base + label_padding),
            label,
            fill="black",
            font=font,
        )

        y_thumbnail = y_base + label_padding + font_size + label_padding

        with Image.open(img_path) as img:
            orig_w, orig_h = img.size

            if placeholder_regions and (start_slide_num + i) in placeholder_regions:
                if img.mode != "RGBA":
                    img = img.convert("RGBA")

                regions = placeholder_regions[start_slide_num + i]

                if slide_dimensions:
                    slide_width_inches, slide_height_inches = slide_dimensions
                else:
                    slide_width_inches = orig_w / CONVERSION_DPI
                    slide_height_inches = orig_h / CONVERSION_DPI

                x_scale = orig_w / slide_width_inches
                y_scale = orig_h / slide_height_inches

                overlay = Image.new("RGBA", img.size, (255, 255, 255, 0))
                overlay_draw = ImageDraw.Draw(overlay)

                for region in regions:
                    px_left = int(region["left"] * x_scale)
                    px_top = int(region["top"] * y_scale)
                    px_width = int(region["width"] * x_scale)
                    px_height = int(region["height"] * y_scale)

                    stroke_width = max(5, min(orig_w, orig_h) // 150)
                    overlay_draw.rectangle(
                        [(px_left, px_top), (px_left + px_width, px_top + px_height)],
                        outline=(255, 0, 0, 255),
                        width=stroke_width,
                    )

                img = Image.alpha_composite(img, overlay)
                img = img.convert("RGB")

            img.thumbnail((width, height), Image.Resampling.LANCZOS)
            w, h = img.size
            tx = x + (width - w) // 2
            ty = y_thumbnail + (height - h) // 2
            grid.paste(img, (tx, ty))

            if BORDER_WIDTH > 0:
                draw.rectangle(
                    [
                        (tx - BORDER_WIDTH, ty - BORDER_WIDTH),
                        (tx + w + BORDER_WIDTH - 1, ty + h + BORDER_WIDTH - 1),
                    ],
                    outline="gray",
                    width=BORDER_WIDTH,
                )

    return grid


if __name__ == "__main__":
    main()
